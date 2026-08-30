import logging
import os
from pathlib import Path
from typing import List, Dict, Any, Optional

import numpy as np
from pypdf import PdfReader
import pdfplumber
from pptx import Presentation
from docx import Document as DocxDocument
import textract

from src.config import settings
from src.database import get_db_session
from src.embeddings.provider import get_embedding_provider
from src.models.entities import Course, Document
from src.vectorstore.pinecone_store import get_vector_store

logger = logging.getLogger(__name__)

try:
    import pypdfium2 as pdfium
except ImportError:
    pdfium = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

class DocumentIngestor:
    """
    Handles the pipeline of reading course materials,
    parsing them into text, chunking, embedding, and indexing.
    """

    def __init__(self):
        self.embedding_provider = get_embedding_provider()
        self.vector_store = get_vector_store()

    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """Extracts text from PDF files using pypdf, pdfplumber, then OCR."""
        text = ""
        # Try pypdf first
        try:
            reader = PdfReader(file_path)
            for page in reader.pages:
                content = page.extract_text()
                if content:
                    text += content + "\n"
        except Exception as e:
            logger.warning(f"pypdf failed for {file_path}: {e}")

        # If pypdf extracted little or no text, try pdfplumber
        if not text.strip():
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        content = page.extract_text()
                        if content:
                            text += content + "\n"
                logger.info(f"Successfully extracted text using pdfplumber for {file_path}")
            except Exception as e:
                logger.error(f"pdfplumber also failed for {file_path}: {e}")

        # If the PDF is scanned, render pages and OCR them.
        if not text.strip():
            text = self._extract_text_from_pdf_with_ocr(file_path)
        return text

    def _extract_text_from_pdf_with_ocr(self, file_path: Path) -> str:
        """OCR fallback for scanned PDFs."""
        if pdfium is None:
            logger.warning("pypdfium2 is not installed; skipping OCR fallback.")
            return ""
        if pytesseract is None:
            logger.warning("pytesseract is not installed; skipping OCR fallback.")
            return ""

        text = ""
        try:
            pdf = pdfium.PdfDocument(str(file_path))
            for page_index in range(len(pdf)):
                page = pdf[page_index]
                bitmap = page.render(scale=2).to_pil()
                content = pytesseract.image_to_string(bitmap)
                if content:
                    text += content + "\n"
                page.close()
            pdf.close()
            if text.strip():
                logger.info(f"Successfully extracted text using OCR for {file_path}")
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {e}")
        return text

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extracts text from PowerPoint files."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {e}")
        return text

    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extracts text from Word DOCX files."""
        text = ""
        try:
            doc = DocxDocument(file_path)
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text += cell.text + "\n"
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}")
        return text

    def _extract_text_from_doc(self, file_path: Path) -> str:
        """Extracts text from legacy Word DOC files using textract."""
        text = ""
        try:
            # textract requires bytes, returns bytes
            extracted = textract.process(str(file_path))
            text = extracted.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"Error extracting DOC {file_path}: {e}")
        return text

    def _chunk_text(self, text: str) -> List[str]:
        """Splits text into overlapping chunks."""
        text = text.replace("\r", "").replace("\n\n", "\n")
        chunks = []
        start = 0
        while start < len(text):
            end = start + settings.CHUNK_SIZE
            chunk = text[start:end]
            chunks.append(chunk)
            start += settings.CHUNK_SIZE - settings.CHUNK_OVERLAP
        return chunks

    def process_course(self, course_code: str, course_title: str, description: Optional[str] = None):
        """
        Processes all documents for a given course.
        Expects files to be in COURSES_DIR / course_code /
        """
        course_dir = settings.COURSES_DIR / course_code
        stats = {"processed": 0, "indexed": 0, "skipped": 0}
        if not course_dir.exists() or not course_dir.is_dir():
            logger.error(f"Course directory not found: {course_dir}")
            return stats

        with get_db_session() as session:
            # 1. Ensure course exists in DB
            course = session.query(Course).filter(Course.code == course_code).first()
            if not course:
                course = Course(code=course_code, title=course_title, description=description)
                session.add(course)
                session.flush()
                logger.info(f"Created course: {course_code}")
            course_id = course.id

        # 2. Process documents outside a DB session. Parsing, model loading,
        # embedding, and vector uploads can take long enough for remote DB
        # providers to close idle connections.
        for file_path in course_dir.iterdir():
            if file_path.is_dir():
                continue

            ext = file_path.suffix.lower()
            if ext not in [".pdf", ".pptx", ".docx", ".doc"]:
                continue

            stats["processed"] += 1
            logger.info(f"Processing document: {file_path.name}")

            # Extract text
            if ext == ".pdf":
                text = self._extract_text_from_pdf(file_path)
            elif ext == ".pptx":
                text = self._extract_text_from_pptx(file_path)
            elif ext == ".docx":
                text = self._extract_text_from_docx(file_path)
            elif ext == ".doc":
                text = self._extract_text_from_doc(file_path)
            else:
                continue

            if not text.strip():
                logger.warning(f"No text extracted from {file_path.name}")
                stats["skipped"] += 1
                continue

            # Chunking
            chunks = self._chunk_text(text)
            if not chunks:
                stats["skipped"] += 1
                continue

            # Embedding
            vectors = self.embedding_provider.embed_documents(chunks)

            # Vector Store Indexing
            ids = [f"{course_code}_{file_path.stem}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "course": course_code,
                    "source": file_path.name,
                    "text": chunk,
                    "chunk_id": i
                }
                for i, chunk in enumerate(chunks)
            ]

            self.vector_store.upsert_documents(
                ids=ids,
                vectors=vectors,
                metadatas=metadatas,
                namespace=course_code
            )

            # DB Record
            with get_db_session() as session:
                doc = session.query(Document).filter(
                    Document.course_id == course_id,
                    Document.filename == file_path.name
                ).first()

                if not doc:
                    doc = Document(
                        course_id=course_id,
                        filename=file_path.name,
                        file_type=ext[1:],
                        file_path=str(file_path),
                        file_size_bytes=file_path.stat().st_size,
                        chunks_count=len(chunks),
                        is_indexed=True
                    )
                    session.add(doc)
                else:
                    doc.chunks_count = len(chunks)
                    doc.is_indexed = True

            logger.info(f"Indexed {len(chunks)} chunks for {file_path.name}")
            stats["indexed"] += 1

        logger.info(f"Successfully processed course {course_code}")
        return stats
